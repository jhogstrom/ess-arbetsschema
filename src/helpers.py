import glob
import logging
import os
import re
from typing import Any, List, Optional

from pptx import Presentation, presentation
from pptx.dml.color import RGBColor


class FileHelper:
    def __init__(self, logger) -> None:
        self.logger = logger

    def make_filename(self, filename: str, *, dirs: Optional[List[str]] = None) -> str:
        """
        Resolves the given filename to an existing file path or Google Sheet ID.
        If the filename matches the Google Sheet ID pattern, it is returned as-is.
        Otherwise, the function searches for the file in the current directory and in the provided
        list of directories (`dirs`). If the file is found, its path is returned. If multiple matches
        are found, the most recent one (by basename, in reverse order) is returned. If the file cannot
        be found, a FileNotFoundError is raised.
        Args:
            filename (str): The name of the file or Google Sheet ID to resolve.
            dirs (Optional[List[str]], optional): A list of directories to search for the file. Defaults to None.
        Returns:
            str: The resolved file path or Google Sheet ID.
        Raises:
            ValueError: If no filename is provided.
            FileNotFoundError: If the file cannot be found in the specified directories.
        """
        google_sheet_id_pattern = r"^[a-zA-Z0-9-_]{44}$"
        is_google_sheet_id = bool(re.match(google_sheet_id_pattern, filename))

        self.logger.debug(f"Filename is Google Sheet ID: {is_google_sheet_id}")
        if is_google_sheet_id:
            return filename

        if not filename:
            raise ValueError("No filename provided")
        if os.path.exists(filename):
            self.logger.debug(f"File {filename} => {filename}")
            return filename
        dirs = dirs or []
        for d in dirs:
            f = os.path.join(d, filename)
            if os.path.exists(f):
                self.logger.debug(f"File {filename} => {f}")
                return f

        matches = []
        for d in dirs:
            p = os.path.join(d, filename)
            self.logger.debug(f"Searching for {filename} in {p}")
            matches.extend(glob.glob(p))
        if not matches:
            raise FileNotFoundError(
                f"File {filename} not found in {dirs} using pattern {filename}"
            )
        result = max(matches, key=os.path.getmtime)
        # result = max([_ for _ in matches if "~" not in _], key=os.path.getmtime)
        self.logger.debug(f"File {filename} => {result}")
        return result

    def read_pptx_file(self, filename: str) -> presentation.Presentation:
        """
        Read the PowerPoint file and return the Presentation object.

        Args:
            filename (str): Filename of the PowerPoint file to read

        Raises:
            ValueError: Raised if no filename is provided
            FileNotFoundError: Raised if the file is not found

        Returns:
            Presentation: PowerPoint Presentation object
        """
        if not filename:
            raise ValueError("No filename provided")
        if not os.path.exists(filename):
            self.logger.warning(f"Powerpoint-file {filename} not found")
            exit(1)

        self.logger.info(f"Reading file {filename}")

        # Open the PowerPoint file and return it
        return Presentation(filename)


def setup_logger(logger_name: str, level: str | None = None) -> logging.Logger:
    logger = logging.getLogger(logger_name)
    level = level or "INFO"
    logger.setLevel(level)
    ch = logging.StreamHandler()
    ch.setLevel(level)
    formatter = logging.Formatter(
        "%(asctime)s - %(name)s - %(levelname)s - %(message)s"
    )
    ch.setFormatter(formatter)
    logger.addHandler(ch)
    return logger


def get_shape(slide, shape_name: str, logger: logging.Logger) -> Optional[Any]:
    """
    Get the shape by member id or member name by looking at
    the text in the shape. If the id or name is in the text,
    return the shape.

    Args:
        slide: Slide in pptx
        member_id (str): Member id to search for

    Returns:
        Shape: Shape, if found by member id or member name, otherwise None
    """
    for shape in slide.shapes:
        if shape.name == shape_name:
            logger.debug(f"Found shape '{shape_name}'.")
            return shape
    for shape in slide.shapes:
        if any(_ in shape.text for _ in shape_name.split()):
            logger.debug(f"Found shape {shape.name} matching {shape_name}")
            return shape
    logger.debug(f"Did not find shape with name '{shape_name}'")
    return None


def make_shape_name(member_id: int) -> str:
    return f"Member: {member_id}"


def color_boats(
    slide,
    members: List[int],
    color: RGBColor,
    logmsg: str,
    logger: logging.Logger,
    terse: bool = False,
):
    for member in members:
        shape = get_shape(slide, make_shape_name(member), logger)
        if shape:
            try:
                shape.fill.fore_color.rgb = color
            except TypeError:
                logger.error(f"Could not set color on shape {shape.text}")
            text = shape.text.replace("\n", "--")
            if not terse:
                logger.info(f"Member {member} ('{text}') {logmsg}")
            shape.name = f"Member: {member}"
        else:
            if not terse:
                logger.warning(f"Member {member} {logmsg}, but not found on map")
