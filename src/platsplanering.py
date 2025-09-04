import argparse
import datetime
import json
import math
import os
import re
from pprint import pprint
from typing import Any, Dict, Hashable, List, Optional

import pandas as pd
import pyperclip
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Cm, Inches, Pt

from googleapi import get_google_sheet, get_sheet_titles
from helpers import FileHelper, color_boats, get_shape, setup_logger

FILL_COLOR = RGBColor(214, 245, 214)
FILL_COLOR_NOSPOT = RGBColor(255, 230, 230)
FILL_COLOR_ON_LAND = RGBColor(230, 230, 255)
FILL_COLOR_MEMBER_LEFT = RGBColor(255, 153, 255)


def define_colors(filename: Optional[str]) -> Dict[str, RGBColor]:
    result = {
        "reserved": RGBColor(214, 245, 214),
        "declined": RGBColor(255, 230, 230),
        "member_left": RGBColor(255, 153, 255),
        "on_land": RGBColor(230, 230, 255),
        "unknown": RGBColor(255, 255, 255),
    }

    if filename and os.path.exists(filename):
        try:
            with open(filename) as f:
                user_colors = json.load(f)
            try:
                for key, value in user_colors.items():
                    result[key] = RGBColor(*value)
            except ValueError:
                logger.error(f"Could not read colors from file {filename}")
        except json.JSONDecodeError:
            logger.error(f"Could not read colors from file {filename}")
    return result


def parseargs():
    year = datetime.datetime.now().year
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--file", default="*karta*.pptx", help="PowerPoint file to read."
    )
    parser.add_argument(
        "--requests",
        default=f"Anmälningar {year}.xlsx",
        help="Excel file with requests for spots.",
    )
    parser.add_argument(
        # "--boats", default="Båtmått*.xlsx", help="Excel file with boat information"
        "--members",
        default="Alla_medlemmar_inkl_båtinfo_*.xlsx",
        help="Excel file with boat information.",
    )
    parser.add_argument(
        "--outfile",
        default=f"stage/varvskarta {year}.pptx",
        help="Filename for output PowerPoint file.",
    )
    parser.add_argument(
        "--exmembers",
        default="boatinfo/ex-members.txt",
        help="Filename with ex-members",
    )
    parser.add_argument(
        "--onland",
        default="boatinfo/sommarliggare.xlsx",
        help="Excel file with members already on land.",
    )
    parser.add_argument(
        "--scheduled",
        default="Torrsättning*.xlsx",
        help="Report from BAS with scheduled members.",
    )
    parser.add_argument("--updateboat", required=False, help="Update boat information")
    return parser.parse_args()


def make_items_integer(thelist: list) -> List[int]:
    """
    Make all items in the list integers.
    If the item is not an integer,
    try to convert it to an integer by
    removing all non-numeric characters.

    The result is deduplicated.

    Args:
        thelist (list): List of items to convert to integers

    Returns:
        List[int]: Returns a list of integers
    """
    logger.debug(f"Cleaning a list of {len(thelist)} 'integers'.")
    for i, m in enumerate(thelist):
        org = m
        if not isinstance(m, int):
            for c in m:
                if not c.isdigit():
                    m = m.replace(c, "")
            try:
                thelist[i] = int(m)
            except ValueError:
                logger.error(f"Could not convert '{org}' to integer")
    result = list(set(thelist))
    logger.info(f"Read {len(result)} valid requests")
    return sorted(result)


COL_TITLE_memberid = "Medlemsnummer"


def read_members_on_land(document: str) -> List[int]:
    year = datetime.datetime.now().year
    if os.path.exists(document):
        logger.info(f"Reading who's on land from file {document}")
        members = pd.read_excel(document)
        return members.loc[members["År"] == year, COL_TITLE_memberid].tolist()

    values = get_google_sheet(document, get_sheet_titles(document)[0])
    headers = values[0] if values else []
    COL_memberid = (
        headers.index(COL_TITLE_memberid) if COL_TITLE_memberid in headers else -1
    )
    COL_year = headers.index(str(year)) if str(year) in headers else -1

    result = sorted(
        [
            int(row[COL_memberid])
            for row in values[1:]
            if row and (row[COL_year] or " ") in "xX"
        ]
    )
    return result


def read_ex_members(document: str) -> list[int]:
    if os.path.exists(document):
        with open(document) as f:
            lines = [_ for _ in f.readlines() if not _.startswith("#")]

        # Regular expression to find the first number in each line
        number_pattern = re.compile(r"\d+")

        # Extract the first number from each line
        return [
            int(number_pattern.search(line).group())
            for line in lines
            if number_pattern.search(line)
        ]

    values = get_google_sheet(document, get_sheet_titles(document)[0])
    headers = values[0] if values else []
    COL_memberid = (
        headers.index(COL_TITLE_memberid) if COL_TITLE_memberid in headers else -1
    )
    result = sorted(
        [
            int(row[COL_memberid])
            for row in values[1:]
            if row and len(row) > COL_memberid and row[COL_memberid]
        ]
    )

    return result


def read_schedule(document: str) -> list[int]:
    logger.info(f"Reading who's scheduled for dry dock from file {document}.")
    schedule = pd.read_excel(document)
    scheduled = sorted(
        [int(_) for _ in set(schedule["Medlemsnr"].tolist()) if not math.isnan(_)]
    )
    logger.info(f"Read {len(scheduled)} scheduled for dry dock from {document}.")
    return scheduled


def get_no_spot_requested(document: str) -> list[int]:
    # TODO: Parametrize the NO_SPOT_OPTION string
    NO_SPOT_OPTION = (
        "Jag vill INTE ta upp min båt i år och vill INTE ha nån vinterplats hos ESS"
    )
    COL_TITLE_upptagning = "Upptagning"
    if os.path.exists(document):
        request_data = pd.read_excel(document)
        result = make_items_integer(
            request_data.loc[
                request_data[COL_TITLE_upptagning] == NO_SPOT_OPTION, COL_TITLE_memberid
            ].tolist()
        )
    else:
        values = get_google_sheet(document, get_sheet_titles(document)[0])
        headers = values[0] if values else []
        COL_memberid = (
            headers.index(COL_TITLE_memberid) if COL_TITLE_memberid in headers else -1
        )
        COL_upptagning = (
            headers.index(COL_TITLE_upptagning)
            if COL_TITLE_upptagning in headers
            else -1
        )
        result = [
            row[COL_memberid]
            for row in values[1:]
            if row[COL_upptagning] == NO_SPOT_OPTION
        ]
        result = sorted(make_items_integer(result))
        logger.info(f"Read {len(result)} who do not want a spot: {result}")

    return result


def read_requests(document: str) -> List[int]:
    """
    Retrieves a list of member IDs from a given file.
    If the file exists locally, it reads the member IDs from the specified column in an Excel file.
    If the file does not exist, it downloads the data from a Google Sheet and extracts the member IDs.
    Args:
        filename (str): The path to the local file or the identifier for the Google Sheet.
    Returns:
        List[int]: A list of member IDs as integers.
    """
    # TODO: Handle members with more than one boat.
    # The basic idea is to use the field containing boatname, and
    # add ".<n>" to the member ID. That would require a type change from integer
    # to float OR use strings and add ".<boatname>" to the member ID.
    if os.path.exists(document):
        logger.info(f"Reading file {document}")
        request_data = pd.read_excel(document)
        requests = request_data[COL_TITLE_memberid].tolist()
    else:
        values = get_google_sheet(document, get_sheet_titles(document)[0])
        headers = values[0] if values else []
        COL_memberid = (
            headers.index(COL_TITLE_memberid) if COL_TITLE_memberid in headers else -1
        )
        requests = [
            row[COL_memberid] for row in values[1:] if row and len(row) > COL_memberid
        ]
        # pprint(values)
        pprint(len(sorted(requests)))
        # exit(1)
    return sorted(make_items_integer(requests))


def read_members(
    document: str, columns: Optional[List[str]] = None
) -> List[Dict[Hashable, Any]]:
    logger.info(f"Reading member file {document}.")
    values = pd.read_excel(document)
    # Filter out the columns we are interested in
    columns = columns or [
        "Medlemsnr",
        "Längd (båt)",
        "Bredd",
        "Förnamn",
        "Efternamn",
        "Plats",
        # "Modell",
    ]

    result = values[columns].to_dict(orient="records")
    logger.info(f"Read {len(result)} boats from member file.")
    return result


def get_boats(
    *,
    members: List[Dict[Hashable, Any]],
    already_there: List[int],
    scheduled: List[int],
    no_spot_requested: List[int],
    requested_spots: List[int],
) -> list:
    # TODO: Parametrize the column names

    logger.info(f"Requested spots: {len(requested_spots)}")
    logger.info(f"Requested spots: {requested_spots}")

    logger.info(f"Booked spots: {len(scheduled)}")
    logger.info(f"Already there: {len(already_there)}")
    logger.info(f"No spot requested: {len(no_spot_requested)}")

    booked_but_not_requested = 0
    for id in scheduled:
        if id not in requested_spots:
            logger.warning(f"Member {id} not in requests, but booked for dry dock.")
            requested_spots.append(id)
            booked_but_not_requested += 1
    if booked_but_not_requested:
        logger.info(f"Added {booked_but_not_requested} scheduled boats to requests.")

    for id in already_there:
        if id not in requested_spots:
            logger.warning(f"Member {id} not in requests, but already on land.")
            requested_spots.append(id)

    for r in requested_spots:
        if r not in [b["Medlemsnr"] for b in members]:
            logger.warning(f"Member {r} not found in member list.")

    requests: List[dict] = [b for b in members if b["Medlemsnr"] in requested_spots]
    for member in requests:
        member["member"] = int(member.pop("Medlemsnr"))
        member["length"] = float(member.pop("Längd (båt)").replace(",", ".")) + 1
        # add 1m to width
        w = float(member.pop("Bredd").replace(",", ".")) + 1
        # Round width up to nearest .0 or .5
        w = math.ceil(w * 2) / 2
        member["width"] = w
        # boat['name'] = f"{boat.pop('Förnamn')[0]} {boat.pop('Efternamn')}\n({boat.pop('Modell')})"
        member["name"] = f"{member.pop('Efternamn')}"
        member["requested"] = member["member"] not in no_spot_requested

    # Make boats unique by member id
    result = list({boat["member"]: boat for boat in requests}.values())
    logger.info(f"After deduplication: {len(result)} unique boats to go on land.")

    return result


def set_shape_text(shape, text: str) -> None:
    """
    Set the text of the shape to the text provided and format it.

    Args:
        shape: Shape in pptx
        text (str): Text to set in the shape
    """
    TEXT_COLOR = RGBColor(0, 0, 0)  # Set the text color to black
    TEXT_MARGIN = Cm(0.1)
    FONT_SIZE = 8
    text_frame = shape.text_frame
    text_frame.margin_left = TEXT_MARGIN
    text_frame.margin_right = TEXT_MARGIN
    text_frame.margin_top = TEXT_MARGIN
    text_frame.margin_bottom = TEXT_MARGIN
    text_frame.clear()  # Clear the existing text
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Make the text
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(FONT_SIZE)
    run.font.color.rgb = TEXT_COLOR
    run.font.bold = False
    # Enable text autofit
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def format_shape(shape, fill_color: RGBColor) -> None:
    """
    Format the shape with the fill color provided.

    Args:
        shape: Shape in pptx
        fill_color (RGBColor): Fill color to set in the shape
    """
    LINE_COLOR = RGBColor(0, 0, 0)  # Black color
    LINE_WIDTH = Inches(0.01)  # Thin outline
    # Set the shape fill color to the color provided
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color

    # Set the shape outline to thin black
    line = shape.line
    line.color.rgb = LINE_COLOR
    line.width = LINE_WIDTH


def ensured_shape(slide, shape_name: str, boat: Dict[str, Any]):
    # TODO: Parametrize the SCALE_LENGTH and SCALE_WIDTH
    SCALE_LENGTH = 5
    SCALE_WIDTH = 5
    # If creating a new shape, set the left and top position to top left on the slide.
    left = Pt(1)
    top = Pt(1)

    shape = get_shape(slide, shape_name, logger)
    if not shape:
        width = Pt(boat["length"] * SCALE_WIDTH)
        length = Pt(boat["width"] * SCALE_LENGTH)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, length)
    else:
        shape.width = Pt(boat["length"] * SCALE_WIDTH)
        shape.height = Pt(boat["width"] * SCALE_LENGTH)
    shape.name = shape_name
    return shape


def add_boats_to_map(
    slide: Slide, boats: list, *, already_there: List[int], ex_members: List[int]
):
    fills = {True: colors["reserved"], False: colors["declined"]}

    logger.info(f"Adding {len(boats)} boats to map")
    counts = {True: 0, False: 0}
    # Iterate over the list and create or reuse shapes
    for boat in boats:
        member = boat["member"]
        shape_name = f"Member: {member}"
        logger.debug(f"Adding boat {shape_name}")
        shape = ensured_shape(slide, shape_name, boat)

        format_shape(shape, fills[boat["requested"]])
        counts[boat["requested"]] += 1

        expansions = {
            "length": boat["length"],
            "width": boat["width"],
            "name": boat["name"],
            "member": member,
        }
        NAME_ONLY = "{member} {name}"  # noqa: F841
        SIZE_ONLY = "{member} {length:.1f}x{width:.1f}"  # noqa: F841
        FULL = "{member} {name}\n{length:.1f}x{width:.1f}"  # noqa: F841
        caption = FULL.format(**expansions)

        set_shape_text(shape, caption)
    logger.info(f"Spots count: {counts[True]}")
    logger.info(f"Yield count: {counts[False]}")
    color_boats(
        slide, list(set(ex_members)), colors["member_left"], "has left the club", logger
    )
    color_boats(slide, already_there, colors["on_land"], "is already on land", logger)

    for s in slide.shapes:
        if "Rectangle" in s.name:
            logger.info(f"'{s.text}' has not requested a spot")
            s.fill.fore_color.rgb = colors["unknown"]
    for boat in boats:
        if not boat["requested"]:
            logger.info(f"{boat['member']} {boat['name']} has declined a spot")


def mark_all_boats_as_unhandled(slide: Slide):
    for s in slide.shapes:
        if "Member:" in s.name:
            s.fill.fore_color.rgb = colors["unknown"]


def remove_shape_by_name(slide, shape_name: str) -> bool:
    """
    Remove the shape by name from the slide object and return True if removed, otherwise False.

    Args:
        slide: Slide in pptx
        shape_name (str): Name of the shape to remove

    Returns:
        bool: True if shape removed, otherwise False
    """
    shape = get_shape(slide, shape_name, logger)
    if shape:
        sp = shape.element
        sp.getparent().remove(sp)
        return True
    return False


def update_revision(shape, revision: str = "1", boats: list | None = None):
    """
    Update the revision shape with the current date and time.

    Args:
        shape: Shape in pptx
    """
    text = [
        f"Revision {revision}",
        f"Båtar: {len(boats)}",
        datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
    ]
    shape.text = "\n".join(text)


def update_title(shape, title: str):
    """
    Update the title shape with the given title.

    Args:
        shape: Shape in pptx
        title (str): Title text to set
    """
    shape.text = title
    shape.text_frame.paragraphs[0].font.size = Pt(35)


def update_legend(slide, colors):
    """
    Update the legend boxes with the current colors.

    Args:
        slide: Slide in pptx
        colors: Dictionary with colors for the legend boxes
    """
    for key, color in colors.items():
        shape_name = f"Legend: {key}"
        shape = get_shape(slide, shape_name, logger)
        if not shape:
            logger.warning(f"Could not find shape '{shape_name}'")
        else:
            logger.debug(f"Setting color {color} for shape '{shape_name}'")
            shape.fill.fore_color.rgb = color


def read_and_process_input(
    *,
    request_source: str,
    members_source: str,
    on_land_source: str,
    scheduled_source: str,
    updateboat: int | None = None,
) -> List[Dict[Hashable, Any]]:

    requested_spots = read_requests(request_source)
    members = read_members(members_source)
    if updateboat:
        scheduled = [int(updateboat)]
        no_spot_requested = []
        already_there = []
    else:
        already_there = read_members_on_land(on_land_source)
        scheduled = read_schedule(scheduled_source)
        no_spot_requested = get_no_spot_requested(request_source)

    boats = get_boats(
        members=members,
        already_there=already_there,
        scheduled=scheduled,
        no_spot_requested=no_spot_requested,
        requested_spots=requested_spots,
    )
    return boats


def send_reminders(
    *, memberfile: str, request_source: str, on_land_source: str, ex_members: List[int]
):
    """
    Send reminders to members about their boat spots.
    """
    logger.info("Sending reminders to members...")
    # Implementation for sending reminders goes here
    members = read_members(
        memberfile,
        columns=[
            "Medlemsnr",
            "Förnamn",
            "Efternamn",
            "Epost 1",
            "Längd (båt)",
            "Bredd",
        ],
    )

    # Remove members without a boat
    members = [
        m for m in members if not pd.isna(m["Längd (båt)"]) and not pd.isna(m["Bredd"])
    ]

    # Skip the opnes who submitted a request
    requests = read_requests(request_source)
    # Skip the ones that are already on land
    already_there = read_members_on_land(on_land_source)
    # Skip the ones that are ex-members
    # ex_members (already passed in as argument)

    handled = set(requests + already_there + ex_members)

    missing = [m for m in members if m["Medlemsnr"] not in handled]
    logger.info(f"Found {len(missing)} members who have not requested a spot.")

    for member in missing:
        logger.info(
            f"Missing {member['Medlemsnr']} {member['Förnamn']} {member['Efternamn']} ({member['Längd (båt)']} x {member['Bredd']})."
        )

    emails = [member["Epost 1"] for member in missing if "Epost 1" in member]
    logger.info(f"Found {len(emails)} email addresses for missing members.")

    # Copy emails to clipboard
    try:
        pyperclip.copy(",".join(emails))
        logger.info(f"Copied {len(emails)} email addresses to clipboard")
    except Exception as e:
        logger.error(f"Failed to copy emails to clipboard: {e}")
        logger.info(f"Email addresses: {emails}")


if __name__ == "__main__":
    args = parseargs()
    colors = define_colors("templates/colors.json")
    logger = setup_logger("spots", "INFO")

    fh = FileHelper(logger)
    members_source = fh.make_filename(args.members, dirs=["boatinfo"])
    ex_members_source = fh.make_filename(args.exmembers, dirs=["boatinfo"])
    on_land_source = fh.make_filename(args.onland, dirs=["boatinfo"])
    request_source = fh.make_filename(args.requests, dirs=["boatinfo"])
    boats = read_and_process_input(
        request_source=request_source,
        members_source=members_source,
        on_land_source=on_land_source,
        scheduled_source=fh.make_filename(args.scheduled, dirs=["boatinfo"]),
        updateboat=args.updateboat,
    )

    ex_members = read_ex_members(ex_members_source)
    already_there = read_members_on_land(on_land_source)

    if args.updateboat:
        logger.info(f"Filtering boats on {args.updateboat}")
        boats = [b for b in boats if b["member"] == int(args.updateboat)]

    powerpoint_filename = fh.make_filename(args.file, dirs=["stage", "templates"])
    logger.info(f"PowerPoint file path: {powerpoint_filename}")

    ppt = fh.read_pptx_file(powerpoint_filename)
    map_slide = ppt.slides[0]
    mark_all_boats_as_unhandled(slide=map_slide)
    add_boats_to_map(
        slide=map_slide, boats=boats, ex_members=ex_members, already_there=already_there
    )
    update_revision(get_shape(map_slide, "Revision", logger), revision="1", boats=boats)
    year = datetime.datetime.now().year
    update_title(
        get_shape(map_slide, "Rubrik", logger), title=f"Varvskarta ESS {year}/{year+1}"
    )
    update_legend(map_slide, colors)
    try:
        ppt.save(args.outfile)
        logger.info(f"Saved file '{args.outfile}'")
    except PermissionError:
        logger.error(f"Could not save file '{args.outfile}'")
        logger.error("File is open in another application")
        exit(1)

    send_reminders(
        memberfile=members_source,
        request_source=request_source,
        on_land_source=on_land_source,
        ex_members=ex_members,
    )
